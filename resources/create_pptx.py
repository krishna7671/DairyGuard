#!/usr/bin/env python3
"""
Convert DairyGuard presentation content to PPTX format
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_dairyguard_pptx():
    # Create presentation object
    prs = Presentation()
    
    # Define colors (Professional Blue #8BAFD0)
    primary_color = RGBColor(139, 175, 208)  # #8BAFD0
    bg_color = RGBColor(18, 18, 18)  # #121212 (dark background)
    text_color = RGBColor(232, 232, 232)  # #E8E8E8 (light text)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "DairyGuard"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = "Real-Time Milk Shelf Life Monitoring\n& Quality Analytics Platform"
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
    
    # Slide 2: Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Value Proposition:
• IoT & ML platform for milk shelf life prediction
• Real-time monitoring of temperature, pH, bacteria, humidity
• 6 user pages: Dashboard, Sensors, Prediction, QC Charts, Alerts, Analytics

Problem Addressed:
• $338B annual food waste crisis (2023)
• 27% of food supply becomes waste
• 80% of surplus comes from perishables

Key Benefits:
• Reduced spoilage and waste
• Lower recall risk through real-time monitoring
• Operational efficiency gains
• Regulatory confidence through SPC-driven control"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 3: Industry Problems - Waste Statistics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Industry Problem: Food Waste Crisis"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Key Statistics (2023):
• Total food supply: 237 million tons
• Unsold/uneaten share: 31% of supply
• Food waste to disposal: 27% (≈63 million tons)
• Economic value of waste: $338 billion
• Perishables dominate: 80% of surplus

Dairy Industry Impact:
• Included in perishables category
• Highly exposed to spoilage and temperature abuse
• Cold chain variability increases risk
• Date label confusion leads to premature discards

Source: ReFED 2024/2025"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 4: Dairy Industry Challenges
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Four Critical Dairy Challenges"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """1. Spoilage and Waste
   • Temperature excursions reduce shelf life
   • Acidification trends go undetected
   • Inventory rotation inefficiencies

2. Process Inconsistency & Quality Defects
   • Unstable processes produce defects
   • Lack of statistical process control
   • No real-time deviation alerts

3. Safety/Regulatory Pressure
   • FDA suspended milk quality testing program (2025)
   • Increased reliance on in-house monitoring
   • Compliance documentation gaps

4. Cold Chain Fragility
   • Distribution variability magnifies risk
   • Limited visibility across supply chain
   • Reactive rather than proactive management"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 5: Solution Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "DairyGuard Solution Overview"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Integrated Technology Platform:

IoT Sensor Integration
• Real-time monitoring of 4 critical parameters
• Temperature, pH, bacterial activity, humidity
• Edge-to-cloud data transmission
• Automated deviation alerts

Machine Learning Predictions
• 92% prediction accuracy
• Confidence intervals for decision support
• Shelf life forecasting with risk analysis
• Real-time quality scoring

Statistical Process Control
• 7 SPC chart types for quality monitoring
• Pareto, X-bar/R, Fishbone analysis
• Process stability detection
• Automated corrective action guidance

6 User Interface Pages
• Dashboard, Sensors, Prediction, QC Charts, Alerts, Analytics"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 6: Real-Time IoT Monitoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Real-Time IoT Monitoring"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Sensor Types & Capabilities:

Temperature Monitoring
• Range: -2°C to 8°C (refrigerated dairy range)
• Precision: ±0.1°C accuracy
• Real-time alerts for excursions

pH Level Tracking
• Range: 6.0 to 7.0 (dairy optimal range)
• Detection of acidification trends
• Early spoilage indicators

Bacterial Activity Analysis
• Real-time bacterial count monitoring
• Subclinical detection capabilities
• Contamination risk assessment

Humidity Control
• Optimal range: 85-95% for dairy storage
• Mold and contamination prevention
• Environmental stability monitoring

Edge-to-Cloud Architecture
• Continuous data transmission
• Local processing capabilities
• Secure data aggregation"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 7: ML Shelf Life Prediction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "ML Shelf Life Prediction System"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Advanced Predictive Analytics:

Core Capabilities
• 92% prediction accuracy with confidence intervals
• Real-time shelf life forecasting
• Multi-parameter risk assessment
• Confidence bounds for decision support

Prediction Models
• Temperature-time relationships
• pH impact on shelf stability
• Bacterial growth modeling
• Environmental factor integration

Decision Support Features
• Risk scoring (Low/Medium/High)
• Confidence intervals (± days)
• Intervention recommendations
• Inventory optimization guidance

Business Impact
• Earlier spoilage detection
• Smarter inventory rotation
• Reduced waste and recalls
• Optimized production scheduling"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 8: Quality Control Charts
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Statistical Process Control (SPC) Charts"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """7 Quality Control Chart Types:

1. Pareto Chart
   • Prioritizes vital few defects
   • 80/20 rule application
   • Root cause prioritization

2. X-bar/R Charts
   • Process mean and range monitoring
   • Statistical process control
   • Out-of-control detection

3. Fishbone Diagram
   • Root cause analysis framework
   • 6M methodology (Man, Machine, Method, Material, Mother Nature, Measurement)

4. Histogram
   • Distribution pattern analysis
   • Process capability assessment
   • Normal vs. abnormal patterns

5. Scatter Plot
   • Correlation analysis
   • Variable relationship identification
   • Trend line fitting

6. P-Chart
   • Proportion defect monitoring
   • Statistical quality control
   • Process stability tracking

7. C-Chart
   • Count-based defect monitoring
   • Constant sample size analysis
   • Quality trend identification"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = text_color
    
    # Slide 9: Application Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Application Features - 6 Core Pages"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Dashboard Page
• Real-time KPIs and alerts overview
• Critical temperature and pH monitoring
• System health indicators

Sensors Page
• Live sensor data visualization
• Individual sensor status monitoring
• Historical data trends

Prediction Page
• ML shelf life predictions with confidence intervals
• Risk scoring and intervention alerts
• Batch-specific predictions

QC Charts Page
• Interactive SPC chart generation
• Statistical analysis tools
• Quality trend monitoring

Alerts Page
• Severity-based alert management
• Automated notification system
• Alert response tracking

Analytics Page
• Comprehensive data analysis
• Performance metrics and trends
• Predictive insights dashboard

Add Product Page
• New milk batch entry and tracking
• Product-specific prediction models
• Batch lifecycle management"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 10: Business Case & ROI
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Business Case & ROI"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Value Drivers:

Cost Reduction
• Reduced waste: 15-25% reduction in spoilage
• Lower energy costs through optimized cooling
• Reduced rework and quality issues

Risk Mitigation
• Recall prevention through early detection
• Regulatory compliance confidence
• Audit-ready documentation

Operational Efficiency
• Real-time monitoring reduces manual checks
• Automated alerts improve response time
• Data-driven decision making

Revenue Protection
• Consistent product quality
• Extended shelf life management
• Improved customer satisfaction

ROI Projections
• 62% first-year returns (industry average)
• Payback period: 12-18 months
• 3-year ROI: 150-250%"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 11: Implementation Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """4-Phase Deployment Plan:

Phase 1: Foundation (Weeks 1-4)
• Sensor installation and calibration
• Network infrastructure setup
• Initial data collection and baseline establishment

Phase 2: System Integration (Weeks 5-8)
• ML model training with facility data
• User interface deployment
• Staff training and onboarding

Phase 3: Quality Control Setup (Weeks 9-12)
• SPC chart configuration
• Alert threshold optimization
• Process validation and testing

Phase 4: Full Production (Weeks 13-16)
• Complete system activation
• Performance monitoring
• Continuous improvement optimization

Ongoing Support
• 24/7 system monitoring
• Regular model retraining
• Continuous improvement assessments"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 12: Competitive Advantages
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Competitive Advantages"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """5 Key Differentiators:

1. Integrated IoT + ML Approach
   • Real-time monitoring combined with predictive analytics
   • Single platform solution vs. fragmented tools

2. Dairy-Specific Intelligence
   • Purpose-built for dairy industry challenges
   • Optimized for milk shelf life prediction

3. Advanced SPC Implementation
   • 7 comprehensive quality control charts
   • Statistical process control expertise

4. Confidence Interval Predictions
   • 92% accuracy with uncertainty quantification
   • Risk-based decision support

5. Real-Time Intervention Capabilities
   • Immediate alert system
   • Automated corrective action guidance
   • Proactive rather than reactive management"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 13: Security & Governance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Security & Governance"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Comprehensive Security Framework:

Data Protection
• End-to-end encryption for all data transmission
• Secure cloud storage with redundancy
• Regular security audits and compliance checks

Authentication & Access Control
• Multi-factor authentication (MFA)
• Role-based access controls (RBAC)
• Audit trail for all system interactions

Regulatory Compliance
• FDA 21 CFR Part 11 compliance ready
• HACCP integration capabilities
• GFSI benchmark alignment

Data Governance
• Data retention policies
• Backup and disaster recovery
• Incident response procedures

Privacy Protection
• GDPR compliance framework
• Data anonymization capabilities
• User consent management"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 14: Technology Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """IoT Edge to Application Layer:

IoT Sensor Layer
• Temperature, pH, bacteria, humidity sensors
• Edge computing capabilities
• Local data processing and buffering

Communication Layer
• Secure IoT protocols (MQTT, CoAP)
• Edge-to-cloud connectivity
• Redundant communication paths

Cloud Processing Layer
• Real-time data ingestion
• ML model inference and training
• Historical data storage and analysis

Application Layer
• Web-based user interface
• Mobile-responsive design
• API integration capabilities

Data Storage
• Time-series database for sensor data
• ML model artifact storage
• User and configuration management"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 15: Future Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Future Roadmap"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Short-term Enhancements (6 months):
• Advanced analytics dashboard
• Mobile application development
• Integration with existing ERP systems

Medium-term Developments (12 months):
• AI-driven anomaly detection
• Predictive maintenance features
• Supply chain visibility integration
• Multi-location management capabilities

Long-term Vision (24 months):
• Blockchain-based traceability
• Advanced ML models with federated learning
• Industry-wide benchmarking platform
• Regulatory compliance automation

Strategic Partnerships
• Dairy equipment manufacturers
• Quality assurance organizations
• Research institutions for continuous innovation"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 16: Success Metrics & KPIs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Success Metrics & KPIs"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Measurable Outcomes:

Operational Metrics
• Waste reduction: 15-25% target
• Spoilage incidents: 50% reduction
• Response time to alerts: <15 minutes
• System uptime: 99.9% availability

Quality Metrics
• Prediction accuracy: 92% target
• Shelf life extension: 10-15% average
• Process stability improvement: 30%
• Customer complaints: 40% reduction

Financial Metrics
• ROI achievement: 62% first-year
• Cost per unit reduction: 8-12%
• Energy savings: 5-10%
• Labor efficiency: 15-20% improvement

Compliance Metrics
• Audit readiness: 95% score
• Documentation completeness: 100%
• Regulatory compliance: Zero violations
• Traceability: Complete batch tracking"""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Slide 17: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Next Steps & Call to Action"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content.text = """Ready to Transform Your Dairy Operations?

Immediate Opportunities:
• Pilot program deployment (6-8 weeks)
• Proof of concept implementation
• ROI validation study

Partnership Options:
• Technology licensing agreement
• Joint development partnership
• White-label solution for equipment manufacturers

Contact Information:
• Demo scheduling and technical consultation
• Custom solution design and scoping
• Implementation timeline and cost planning

Value Demonstration:
• Industry benchmark analysis
• Waste reduction potential assessment
• Financial impact projection

Let's discuss how DairyGuard can solve your specific dairy quality challenges and drive measurable business results."""
    
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.font.color.rgb = text_color
    
    # Save the presentation
    output_path = "/workspace/DairyGuard-Business-Presentation.pptx"
    prs.save(output_path)
    print(f"PPTX presentation saved to: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_dairyguard_pptx()